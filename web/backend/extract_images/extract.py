import numpy as np
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import cv2
import numpy as np
import tempfile
from docx import Document
import fitz  # PyMuPDF
from PIL import Image, ImageDraw
from pdf2docx import Converter
from paddleocr import PaddleOCR
import aspose.slides as slides

########################## 1.Preprocess ##########################

def is_single_color(image):
    image_array = np.array(image)
    min_val, max_val = image_array.min(), image_array.max()
    return max_val - min_val < 10

def is_extreme_aspect_ratio(idx, image, threshold=4):
    width, height = image.size
    aspect_ratio = max(height / width, width / height)
    # print(f"idx{idx}: aspect_ratio: {aspect_ratio}")
    return aspect_ratio > threshold

def count_unique_colors(idx, image):
    image_array = np.array(image)
    unique_colors = set(tuple(pixel) for row in image_array for pixel in row)
    # print(f"idx{idx}: unique colors: {len(unique_colors)}")
    return len(unique_colors)

def add_white_background(image):
    if image.mode == 'RGBA':
        # 알파 채널을 처리하여 흰색 배경으로 변환
        background = Image.new('RGBA', image.size, (255, 255, 255, 255))
        background.paste(image, (0, 0), image)
        return background.convert('RGB')  # 최종적으로 RGB로 변환
    else:
        return image.convert('RGB')  # RGBA가 아닌 경우 그냥 RGB로 변환

########################## 2.pdf extract calculate func ##########################

# 이진화된 이미지에서 재귀적으로 문서내의 각 섹션의 경계를 찾아냄
# 주로 텍스트나 그래픽이 포함된 블록을 식별
def find_section_borders(img_binary:np.array, 
                    min_w:float=0.0, min_h:float=0.0, 
                    min_dx:float=15.0, min_dy:float=15.0): 

    def cut_sections(arr:np.array, top_left:tuple, res:list, 
                    min_w:float, min_h:float, min_dx:float, min_dy:float):
        x0, y0 = top_left
        h, w = arr.shape

        projection = np.count_nonzero(arr==255, axis=1)
        pos_y = identify_gaps(projection, min_w, min_dy)
        if not pos_y: return        

        arr_y0, arr_y1 = pos_y
        for r0, r1 in zip(arr_y0, arr_y1):
            x_arr = arr[r0:r1, 0:w]
            projection = np.count_nonzero(x_arr==255, axis=0)
            pos_x = identify_gaps(projection, min_h, min_dx)
            if not pos_x: continue

            arr_x0, arr_x1 = pos_x
            if len(arr_x0)==1:
                res.append((x0+arr_x0[0], y0+r0, x0+arr_x1[0], y0+r1))
                continue
            
            for c0, c1 in zip(arr_x0, arr_x1):
                y_arr = arr[r0:r1, c0:c1]
                top_left = (x0+c0, y0+r0)
                cut_sections(y_arr, top_left, res, min_w, min_h, min_dx, min_dy)

    res = []
    cut_sections(arr=img_binary, top_left=(0, 0), res=res, 
            min_w=min_w, min_h=min_h, min_dx=min_dx, min_dy=min_dy)
    return res

# 최소값과 갭을 기준으로 객체의 시작과 끝을 찾음, 객체의 위치를 찾는데 사용
def identify_gaps(arr_values:np.array, min_value:float, min_gap:float):

    idx = np.where(arr_values>min_value)[0]
    if not len(idx): return

    gaps = idx[1:] - idx[0:-1]
    gap_idx = np.where(gaps>min_gap)[0]
    segment_start = idx[gap_idx]
    segment_end = idx[gap_idx+1]

    starts = np.insert(segment_end, 0, idx[0])
    ends = np.append(segment_start, idx[-1])
    ends += 1 

    return starts, ends

class CorningCustomExtractor:
    def __init__(self, page: fitz.Page) -> None:
        self._page = page

    # 경계 상자에 따라 페이지의 특정 영역을 pixmap 으로 변환
    # pixmap = 문서의 특정 부분을 이미지로 렌더링
    def clip_page_to_pixmap(self, bbox: fitz.Rect = None, zoom: float = 3.0):
        if bbox is None:
            clip_bbox = self._page.rect
        elif self._page.rotation:
            clip_bbox = bbox * self._page.rotation_matrix
        else:
            clip_bbox = bbox
        clip_bbox = self._page.rect & clip_bbox
        matrix = fitz.Matrix(zoom, zoom)
        pix = self._page.get_pixmap(clip=clip_bbox, matrix=matrix)
        return pix

    # svg 윤곽선을 감지하여 각 이미지의 외부 및 내부 경계 분석 
    # 각 페이지 => 흑백 이미지로 변환후 이진화를 통해 경계를 찾아냄
    # 그후 클립하여 저장 
    def detect_svg_contours(self, page_num, min_svg_gap_dx: float, min_svg_gap_dy: float):
        
        pixmap = self.clip_page_to_pixmap(zoom=1.0)
        src = self._pixmap_to_cv_image(pixmap)
        gray = cv2.cvtColor(src, cv2.COLOR_BGR2GRAY)
        _, binary = cv2.threshold(gray, 253, 255, cv2.THRESH_BINARY_INV)
        external_bboxes = find_section_borders(binary, min_dx=min_svg_gap_dx, min_dy=min_svg_gap_dy)

        results = []
        for i, bbox in enumerate(external_bboxes):
            pixmap = self.clip_page_to_pixmap(bbox=bbox, zoom=3.0)
            image = Image.open(io.BytesIO(pixmap.tobytes()))
            if count_unique_colors(i, image) > 300:
                results.append(image)
        
        return results

    @staticmethod
    def _pixmap_to_cv_image(pixmap: fitz.Pixmap):
        img_byte = pixmap.tobytes()
        return cv2.imdecode(np.frombuffer(img_byte, np.uint8), cv2.IMREAD_COLOR)


########################## 3.Extract Images from PPTX, DOCX, PDF ##########################

def extract_images_from_pptx(pptx_bytes):
    output_images = []
    with slides.Presentation(io.BytesIO(pptx_bytes)) as presentation:
        pdf_stream = io.BytesIO()
        presentation.save(pdf_stream, slides.export.SaveFormat.PDF)
        pdf_stream.seek(0)  # Rewind the stream to the beginning
        output_images = extract_images_from_pdf2(pdf_stream.getvalue())

    return output_images
    
    output_images = []
    for picture in iter_picture_shapes(prs):
        image = picture.image.blob
        image = Image.open(io.BytesIO(image))
        output_images.append(image)
    
    return output_images

def extract_images_from_docx(docx_bytes):
    doc = Document(io.BytesIO(docx_bytes))
    
    i = 0
    output_images = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image = rel.target_part.blob
            image = Image.open(io.BytesIO(image))

            if is_extreme_aspect_ratio(i, image):
                continue
                
            # if count_unique_colors(i, image) < 50:
            #     print(f"image{i}.jpg 컬러 스킵")
            #     continue

            image = add_white_background(image)
            output_images.append(image)

    return output_images

#version 1 => 시연용 속도 빠름
def extract_images_from_pdf1(pdf_bytes):
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_pdf:
        temp_pdf.write(pdf_bytes)
        temp_pdf_path = temp_pdf.name

    try:
        pdf_to_docx = io.BytesIO()
        cv = Converter(temp_pdf_path)
        cv.convert(pdf_to_docx, start=0, end=None)
        cv.close()
        pdf_to_docx.seek(0)
        return extract_images_from_docx(pdf_to_docx.getvalue())
    finally:
        os.remove(temp_pdf_path)

#version2 => 제출용, 완벽히 뽑으나 조금 느리고 텍스트도 조금 뽑힘, 직사각형 모양으로 pixel단위를 찾는것임 

def extract_images_from_pdf2(pdf_bytes):
    pdf_file = fitz.open(stream=pdf_bytes, filetype="pdf")
    
    output_images = []
    # PDF 전체 페이지에서 이미지 추출
    for page_number in range(pdf_file.page_count):
        page = pdf_file.load_page(page_number)
        extractor = CorningCustomExtractor(page)
        # print(f"Extracting images from page {page_number + 1}...")
        images = extractor.detect_svg_contours(page_number+1, min_svg_gap_dx=10.0, min_svg_gap_dy=10.0)
        output_images.extend(images)
    
    return output_images
        
########################## Extract text from image ##########################
ocr = PaddleOCR(use_angle_cls=True, lang='en')

def extract_images_from_text(figure_path):
    result = ocr.ocr(figure_path, cls=True)
    text_result = []
    # 인식된 텍스트를 출력합니다.
    for line in result:
        for word in line:
            print(word[1][0])
            text_result.append(word[1][0])
    text_result = '\n'.join(text_result)
    return text_result
