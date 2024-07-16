import cv2
import numpy as np
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import os
import io
import cv2
import numpy as np
from docx import Document
import fitz  # PyMuPDF
from PIL import Image
from pdf2docx import Converter
from paddleocr import PaddleOCR

########################## Preprocess ##########################

def is_single_color(image):
    image_array = np.array(image)
    min_val, max_val = image_array.min(), image_array.max()
    return max_val - min_val < 10

def is_extreme_aspect_ratio(idx, image, threshold=4):
    width, height = image.size
    aspect_ratio = max(height / width, width / height)
    print(f"idx{idx}: aspect_ratio: {aspect_ratio}")
    return aspect_ratio > threshold

def count_unique_colors(idx, image):
    image_array = np.array(image)
    unique_colors = set(tuple(pixel) for row in image_array for pixel in row)
    print(f"idx{idx}: unique colors: {len(unique_colors)}")
    return len(unique_colors)

def add_white_background(image):
    if image.mode == 'RGBA':
        # 알파 채널을 처리하여 흰색 배경으로 변환
        background = Image.new('RGBA', image.size, (255, 255, 255, 255))
        background.paste(image, (0, 0), image)
        return background.convert('RGB')  # 최종적으로 RGB로 변환
    else:
        return image.convert('RGB')  # RGBA가 아닌 경우 그냥 RGB로 변환


########################## Extract Images from PPTX, DOCX, PDF ##########################

def extract_images_from_pptx(pptx_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)

    def iter_picture_shapes(prs):
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    yield shape
    prs = Presentation(pptx_path)
    i = 0
    for picture in iter_picture_shapes(prs):
        image = picture.image
        image_bytes = image.blob
        output_image = Image.open(io.BytesIO(image_bytes))

        if (not is_single_color(output_image)) and (not is_extreme_aspect_ratio(i, output_image)) and (count_unique_colors(i, output_image) > 50): 
            # 투명 배경을 흰색 배경으로 변환
            image = add_white_background(output_image)

            # 이미지 파일을 저장
            output_path = os.path.join(output_folder, f'image{i}.jpg')
            image.save(output_path)
            
            print(f"Saved image{i}.jpg at {output_path}")
            i += 1


def extract_images_from_docx(docx_path, output_folder):

    doc = Document(docx_path)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    i = 0
    # DOCX 문서의 모든 관계(rels)를 순회 처리
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            # 이미지 데이터를 바이트 배열('blob')로 가져옴 
            image_data = rel.target_part.blob
            # 바이트 배열을 Pillow 이미지로 변환
            image = Image.open(io.BytesIO(image_data))
            if is_extreme_aspect_ratio(i, image):
                continue

            # 투명 배경을 흰색 배경으로 변환
            image = add_white_background(image)

            # 이미지 파일을 저장
            output_path = os.path.join(output_folder, f'image{i}.jpg')
            image.save(output_path)
            
            print(f"Saved image{i}.jpg at {output_path}")
            i += 1

#version 1
def extract_images_from_pdf(pdf_path, output_folder):
    # 메모리 내에서 PDF를 DOCX로 변환
    pdf_to_docx = io.BytesIO()
    cv = Converter(pdf_path)
    cv.convert(pdf_to_docx, start=0, end=None)
    cv.close()
    pdf_to_docx.seek(0)  # 메모리 파일의 시작점으로 이동 즉, pdf => word => image로 변환하는 거임 
    extract_images_from_docx(pdf_to_docx, output_folder)

########################## Extract text from image ##########################
ocr = PaddleOCR(use_angle_cls=True, lang='en')

def extract_images_from_text(figure_path):
    result = ocr.ocr(figure_path, cls=True)
    text_result = []
    # 인식된 텍스트를 출력합니다.
    for line in result:
        for word in line:
            print(word[1][0])
            text_result.append(word[1][0])
    text_result = '\n'.join(text_result)
    return text_result
