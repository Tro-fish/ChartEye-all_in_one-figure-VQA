import cv2
import numpy as np
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import os
import io
import cv2
import numpy as np
from docx import Document
from pdf2docx import Converter


##########################preprocess############################
def is_single_color(image):
    min_val, max_val = image.min(), image.max()
    return max_val - min_val < 10

def is_extreme_aspect_ratio(idx, image, threshold=4):
    height, width, _ = image.shape
    aspect_ratio = max(height / width, width / height)
    print(f"idx{idx}: aspect_ratio: {aspect_ratio}")
    return aspect_ratio > threshold

def count_unique_colors(idx, image):
    unique_colors = set(tuple(pixel) for row in image for pixel in row)
    print(f"idx{idx}: unique colors: {len(unique_colors)}")
    return len(unique_colors)

##########################extract##############################

def extract_images_from_pptx(pptx_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)

    def iter_picture_shapes(prs):
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    yield shape
    prs = Presentation(pptx_path)
    i = 0
    for picture in iter_picture_shapes(prs):
        image = picture.image
        image_bytes = image.blob
        output_image = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)

        if (not is_single_color(output_image)) and (not is_extreme_aspect_ratio(i, output_image)) and (count_unique_colors(i, output_image) > 50):
            cv2.imwrite(os.path.join(output_folder, f'{i}.jpg'), output_image)
            i += 1

def extract_images_from_docx(docx_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    doc = Document(docx_path)
    i = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image = rel.target_part.blob
            image_array = np.frombuffer(image, np.uint8)
            image = cv2.imdecode(image_array, cv2.IMREAD_COLOR)

            if is_extreme_aspect_ratio(i, image):
                print(f"Skipped image{i}.jpg 비율 스킵")
                continue
            # if count_unique_colors(i, image) < 50:
            #     print(f"image{i}.jpg 컬러 스킵")
            #     continue
            image_path = os.path.join(output_folder, f'image{i}.jpg')
            with open(image_path, 'wb') as f:
                f.write(cv2.imencode('.jpg', image)[1].tobytes())
                print(f"Saved {image_path}")
            i += 1

def extract_images_from_pdf(pdf_path, output_folder):
    # 메모리 내에서 PDF를 DOCX로 변환
    pdf_to_docx = io.BytesIO()
    cv = Converter(pdf_path)
    cv.convert(pdf_to_docx, start=0, end=None)
    cv.close()
    pdf_to_docx.seek(0)  # 메모리 파일의 시작점으로 이동 즉, pdf => word => image로 변환하는 거임 
    extract_images_from_docx(pdf_to_docx, output_folder)
